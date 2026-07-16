from docx import Document as WordDocument
from openpyxl import Workbook
from pptx import Presentation

from backend.core.document_loader import load_document


def test_load_text_document(tmp_path):
    path = tmp_path / "sample.txt"
    path.write_text("hello knowledge base", encoding="utf-8")
    docs = load_document(str(path))
    assert docs[0].page_content == "hello knowledge base"


def test_load_office_documents(tmp_path):
    word_path = tmp_path / "sample.docx"
    word = WordDocument()
    word.add_paragraph("word content")
    word.save(word_path)

    excel_path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    workbook.active.append(["excel", "content"])
    workbook.save(excel_path)

    slides_path = tmp_path / "sample.pptx"
    slides = Presentation()
    slide = slides.slides.add_slide(slides.slide_layouts[1])
    slide.shapes.title.text = "slide content"
    slides.save(slides_path)

    assert "word content" in load_document(str(word_path))[0].page_content
    assert "excel\tcontent" in load_document(str(excel_path))[0].page_content
    assert "slide content" in load_document(str(slides_path))[0].page_content
