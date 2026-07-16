from pathlib import Path

from docx import Document as WordDocument
from langchain_core.documents import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


SUPPORTED_SUFFIXES = {".txt", ".md", ".docx", ".pdf", ".xlsx", ".pptx"}


def _document(text: str, source: Path, **metadata) -> Document:
    return Document(page_content=text, metadata={"source": str(source), **metadata})


def load_document(file_path: str) -> list[Document]:
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        raise ValueError(f"不支持的文件格式: {suffix}")

    if suffix in {".txt", ".md"}:
        return [_document(path.read_text(encoding="utf-8"), path)]

    if suffix == ".pdf":
        reader = PdfReader(str(path))
        documents = []
        for index, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                documents.append(_document(text, path, page=index))
        return documents

    if suffix == ".docx":
        word = WordDocument(str(path))
        parts = [paragraph.text for paragraph in word.paragraphs if paragraph.text.strip()]
        for table in word.tables:
            parts.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
        return [_document("\n".join(parts), path)]

    if suffix == ".xlsx":
        workbook = load_workbook(path, read_only=True, data_only=True)
        try:
            documents = []
            for sheet in workbook.worksheets:
                rows = [
                    "\t".join("" if value is None else str(value) for value in row)
                    for row in sheet.iter_rows(values_only=True)
                ]
                text = "\n".join(row for row in rows if row.strip())
                if text:
                    documents.append(_document(text, path, sheet=sheet.title))
            return documents
        finally:
            workbook.close()

    presentation = Presentation(str(path))
    documents = []
    for index, slide in enumerate(presentation.slides, 1):
        text = "\n".join(
            shape.text for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        if text:
            documents.append(_document(text, path, slide=index))
    return documents


def load_documents_from_folder(folder_path: str) -> list[Document]:
    documents = []
    for file_path in Path(folder_path).iterdir():
        if file_path.suffix.lower() in SUPPORTED_SUFFIXES:
            documents.extend(load_document(str(file_path)))
    return documents
